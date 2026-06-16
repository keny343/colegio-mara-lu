#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para criar PowerPoint da Apresentação
Sistema de Gestão Escolar - Colégio Mara e Lu
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Cores institucionais
COR_PRINCIPAL = RGBColor(25, 65, 120)  # Azul escuro
COR_SECUNDARIA = RGBColor(76, 175, 80)  # Verde
COR_TEXTO = RGBColor(33, 33, 33)  # Cinzento escuro

def adicionar_titulo_subtitulo(prs, titulo, subtitulo=""):
    """Adiciona slide com título e subtítulo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = titulo
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = COR_PRINCIPAL
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    if subtitulo:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitulo
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = COR_TEXTO
        subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def adicionar_slide_conteudo(prs, titulo, conteudo_list):
    """Adiciona slide com título e conteúdo em bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Fundo do título
    titulo_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    titulo_shape.fill.solid()
    titulo_shape.fill.fore_color.rgb = COR_PRINCIPAL
    titulo_shape.line.color.rgb = COR_PRINCIPAL
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = titulo
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(conteudo_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COR_TEXTO
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

# ========== SLIDE 1: CAPA ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COR_PRINCIPAL

# Título Principal
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.paragraphs[0]
title_p.text = "Implementação de um Sistema de Gestão Escolar"
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_ALIGN.CENTER

# Subtítulo
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.text = "para o Colégio Mara e Lu"
subtitle_p.font.size = Pt(32)
subtitle_p.font.color.rgb = COR_SECUNDARIA
subtitle_p.alignment = PP_ALIGN.CENTER

# Informações
info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
info_frame = info_box.text_frame
info_frame.word_wrap = True

for texto in [
    "Autor: Adnircio do Rosário Quiteculo Inocêncio",
    "Orientador: Plácido Dias",
    "Universidade Metodista de Angola",
    "Junho de 2025"
]:
    p = info_frame.add_paragraph() if info_frame.paragraphs[0].text else info_frame.paragraphs[0]
    if info_frame.paragraphs[0].text:
        p = info_frame.add_paragraph()
    p.text = texto
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 2: PROBLEMA ==========
adicionar_slide_conteudo(prs, "Por que um Sistema?", [
    "❌ Processos 100% manuais (fichas, planilhas)",
    "❌ Perda e duplicidade de dados",
    "❌ Lentidão na triagem de inscrições",
    "❌ Dificuldade no controlo administrativo",
    "✅ Solução: Sistema informatizado integrado",
    "📊 70% das escolas em Angola ainda usam fichas físicas"
])

# ========== SLIDE 3: OBJECTIVOS ==========
adicionar_slide_conteudo(prs, "Objectivos Principais", [
    "📋 Automatizar inscrições e triagem",
    "👥 Centralizar gestão de alunos e turmas",
    "📊 Facilitar lançamento de notas e faltas",
    "📚 Partilhar materiais pedagógicos",
    "🔒 Garantir segurança dos dados",
    "📈 Gerar relatórios administrativos"
])

# ========== SLIDE 4: REQUISITOS FUNCIONAIS ==========
adicionar_slide_conteudo(prs, "Funcionalidades Principais", [
    "✓ Inscrição pública com anexos",
    "✓ Triagem e validação de inscrições",
    "✓ Matrículas automáticas",
    "✓ Gestão de utilizadores (Admin, Coordenador, Professor, Aluno)",
    "✓ Lançamento de notas por disciplina",
    "✓ Registo e justificação de faltas",
    "✓ Gestão de turmas e disciplinas",
    "✓ Upload/download de materiais"
])

# ========== SLIDE 5: ACTORES ==========
adicionar_slide_conteudo(prs, "Actores e Papéis do Sistema", [
    "👨‍💼 Administrador: Gestão global, aprovação de inscrições",
    "📋 Coordenador: Triagem, validação, gestão académica",
    "👨‍🏫 Professor: Notas, faltas, disponibilização de materiais",
    "👨‍🎓 Aluno/Encarregado: Consulta de notas, faltas, materiais",
    "",
    "Fluxo: Candidato → Coordenador (triagem) → Admin (aprovação) → Matrícula"
])

# ========== SLIDE 6: ARQUITECTURA ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Fundo do título
titulo_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
titulo_shape.fill.solid()
titulo_shape.fill.fore_color.rgb = COR_PRINCIPAL
titulo_shape.line.color.rgb = COR_PRINCIPAL

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Arquitectura do Sistema"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)

# Diagrama de blocos
boxes = [
    ("Front-end\n(React SPA)", Inches(3), Inches(1.8), COR_SECUNDARIA),
    ("Back-end\n(Node.js/Express)", Inches(3), Inches(3.5), COR_PRINCIPAL),
    ("Base de Dados\n(MySQL + Storage)", Inches(3), Inches(5.2), RGBColor(200, 50, 50))
]

for label, left, top, color in boxes:
    box = slide.shapes.add_shape(1, left, top, Inches(4), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = RGBColor(100, 100, 100)
    
    text_frame = box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Setas
    if label != "Base de Dados\n(MySQL + Storage)":
        arrow = slide.shapes.add_connector(1, Inches(5), top + Inches(0.9), Inches(5), top + Inches(1.7))
        arrow.line.color.rgb = RGBColor(100, 100, 100)

# Autenticação
auth_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(8.6), Inches(0.8))
auth_frame = auth_box.text_frame
auth_p = auth_frame.paragraphs[0]
auth_p.text = "🔐 Autenticação: JWT com Bcrypt | 🛡️ Middleware para controlo de acesso"
auth_p.font.size = Pt(14)
auth_p.font.color.rgb = COR_TEXTO
auth_p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 7: TECNOLOGIAS ==========
adicionar_slide_conteudo(prs, "Stack Tecnológico", [
    "🎨 Front-end: React, React Router, Axios",
    "   → Design responsivo (Desktop + Mobile)",
    "",
    "⚙️ Back-end: Node.js, Express.js, JWT, Bcrypt, Multer",
    "   → Autenticação segura com hashing",
    "",
    "💾 Banco de Dados: MySQL (relacional)",
    "   → Upload de ficheiros (avatares, materiais, documentos)"
])

# ========== SLIDE 8: FLUXO DE INSCRIÇÃO ==========
adicionar_slide_conteudo(prs, "Fluxo de Inscrição e Matrícula", [
    "1️⃣  Inscrição: Candidato preenche formulário + envia documentos",
    "",
    "2️⃣  Triagem: Coordenador valida inscrição",
    "",
    "3️⃣  Aprovação: Administrador aprova ou rejeita",
    "",
    "4️⃣  Matrícula: Sistema converte em matrícula automática",
    "",
    "5️⃣  Aluno Activo: Aluno entra no sistema com credenciais"
])

# ========== SLIDE 9: GESTÃO ACADÉMICA ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Fundo do título
titulo_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
titulo_shape.fill.solid()
titulo_shape.fill.fore_color.rgb = COR_PRINCIPAL
titulo_shape.line.color.rgb = COR_PRINCIPAL

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Gestão Académica e Pedagógica"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)

# 3 colunas
columns = [
    ("📝 NOTAS", [
        "• Lançamento por disciplina",
        "• Consulta por aluno",
        "• Histórico de períodos"
    ]),
    ("✋ FALTAS", [
        "• Registo por aula",
        "• Justificações",
        "• Relatórios por período"
    ]),
    ("📚 MATERIAIS", [
        "• Upload por professor",
        "• Download por aluno",
        "• Organização por disciplina"
    ])
]

for i, (titulo_col, items) in enumerate(columns):
    left = Inches(0.3 + i * 3.2)
    
    # Título coluna
    col_title = slide.shapes.add_textbox(left, Inches(1.3), Inches(3), Inches(0.5))
    col_tf = col_title.text_frame
    col_p = col_tf.paragraphs[0]
    col_p.text = titulo_col
    col_p.font.size = Pt(18)
    col_p.font.bold = True
    col_p.font.color.rgb = COR_PRINCIPAL
    col_p.alignment = PP_ALIGN.CENTER
    
    # Conteúdo coluna
    col_content = slide.shapes.add_textbox(left, Inches(2), Inches(3), Inches(4.5))
    col_tf = col_content.text_frame
    col_tf.word_wrap = True
    
    for j, item in enumerate(items):
        if j == 0:
            p = col_tf.paragraphs[0]
        else:
            p = col_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = COR_TEXTO
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# ========== SLIDE 10: BASE DE DADOS ==========
adicionar_slide_conteudo(prs, "Modelo de Dados (Entidades)", [
    "👥 Utilizadores: admins, coordenadores, professores, alunos",
    "",
    "📋 Inscrições: inscrições → matrículas (fluxo de aprovação)",
    "",
    "🏫 Académico: turmas, disciplinas, atribuição de professores",
    "",
    "📝 Notas: notas por disciplina e período",
    "",
    "✋ Faltas: faltas por aula, justificações",
    "",
    "🔍 Validações: unicidade, integridade referencial, índices"
])

# ========== SLIDE 11: INTERFACE E UX ==========
adicionar_slide_conteudo(prs, "Interfaces Principais", [
    "👨‍💼 Dashboard Admin:",
    "   • Painel de controlo com estatísticas",
    "   • Gestão de inscrições (pendentes, aprovadas)",
    "",
    "👨‍🏫 Dashboard Professor:",
    "   • Lista de alunos por turma",
    "   • Lançamento de notas e faltas",
    "   • Upload de materiais",
    "",
    "👨‍🎓 Portal Aluno:",
    "   • Perfil e consulta de notas/faltas",
    "   • Download de materiais"
])

# ========== SLIDE 12: RESULTADOS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Fundo do título
titulo_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
titulo_shape.fill.solid()
titulo_shape.fill.fore_color.rgb = COR_PRINCIPAL
titulo_shape.line.color.rgb = COR_PRINCIPAL

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Benefícios e Impacto"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)

# Observados
obs_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.4))
obs_tf = obs_title.text_frame
obs_p = obs_tf.paragraphs[0]
obs_p.text = "Resultados Observados:"
obs_p.font.size = Pt(20)
obs_p.font.bold = True
obs_p.font.color.rgb = COR_PRINCIPAL

obs_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.5))
obs_frame = obs_box.text_frame
obs_frame.word_wrap = True
for texto in ["✅ Melhor organização de inscrições (fluxos centralizados)",
              "✅ Redução do tempo de processamento manual",
              "✅ Maior responsabilização através de controlo de acessos"]:
    if obs_frame.paragraphs[0].text:
        p = obs_frame.add_paragraph()
    else:
        p = obs_frame.paragraphs[0]
    p.text = texto
    p.font.size = Pt(16)
    p.font.color.rgb = COR_TEXTO
    p.space_after = Pt(6)

# Esperados
exp_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.6), Inches(0.4))
exp_tf = exp_title.text_frame
exp_p = exp_tf.paragraphs[0]
exp_p.text = "Resultados Esperados:"
exp_p.font.size = Pt(20)
exp_p.font.bold = True
exp_p.font.color.rgb = COR_PRINCIPAL

exp_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2.8))
exp_frame = exp_box.text_frame
exp_frame.word_wrap = True
for texto in ["📈 Agilidade na triagem e comunicação com encarregados",
              "🔒 Dados seguros, consistentes e organizados",
              "📊 Fácil geração de relatórios administrativos",
              "🚀 Redução de erros e duplicidades",
              "⚡ Eficiência operacional aumentada"]:
    if exp_frame.paragraphs[0].text:
        p = exp_frame.add_paragraph()
    else:
        p = exp_frame.paragraphs[0]
    p.text = texto
    p.font.size = Pt(16)
    p.font.color.rgb = COR_TEXTO
    p.space_after = Pt(6)

# ========== SLIDE 13: CONCLUSÃO ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COR_PRINCIPAL

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Próximos Passos"
title_p.font.size = Pt(48)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_ALIGN.CENTER

content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(8.6), Inches(4.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

recomendacoes = [
    "🔧 Deploy em produção com backups automatizados",
    "⚡ Testes de carga em cenários de pico",
    "📚 Documentação operacional completa",
    "🐳 Containerização com Docker",
    "",
    "Sistema pronto para operação",
    "Arquitetura modular e extensível",
    "Alinha-se com boas práticas internacionais"
]

for i, texto in enumerate(recomendacoes):
    if i == 0:
        p = content_frame.paragraphs[0]
    else:
        p = content_frame.add_paragraph()
    p.text = texto
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(4)
    p.space_after = Pt(4)
    p.alignment = PP_ALIGN.CENTER

# Guardar apresentação
output_path = r"c:\Users\kenyd\Desktop\colegio-mara-lu\APRESENTACAO_SISTEMA_GESTAO_ESCOLAR.pptx"
prs.save(output_path)

print(f"✅ PowerPoint criado com sucesso!")
print(f"📍 Localização: {output_path}")
print(f"📊 Total de slides: {len(prs.slides)}")
